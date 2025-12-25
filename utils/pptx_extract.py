
import io
import os
import numpy as np
from pptx import Presentation
from PIL import Image

def extract_clean_images_from_pptx(pptx_path, output_dir):
    """
    Extract images from PPTX, combining multiple images per slide into a single flattened image.
    Preserves the original layout and positioning of images on the slide.
    """
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    extracted = []
    
    # Get slide dimensions (in EMU - English Metric Units)
    # 1 inch = 914400 EMU
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides, start=1):
        image_data_list = []  # List of (image_array, left, top, width, height) tuples
        
        # Collect all images from this slide with their positions
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.shape_type == 13:  # Skip if not a picture
                continue
            if not hasattr(shape, "image"):
                continue

            image = shape.image
            try:
                img_data = io.BytesIO(image.blob)
                with Image.open(img_data) as im:
                    # Convert to RGB if needed
                    if im.mode != 'RGB':
                        im = im.convert('RGB')
                    img_array = np.array(im)
                    
                    # Get shape position and size (in EMU)
                    left_emu = shape.left
                    top_emu = shape.top
                    width_emu = shape.width
                    height_emu = shape.height
                    
                    image_data_list.append((img_array, left_emu, top_emu, width_emu, height_emu))
            except Exception as e:
                print(f"Skipping image from slide {slide_idx}: {e}")
                continue
        
        if not image_data_list:
            continue
        
        # Process slide images
        if len(image_data_list) == 1:
            # Single image: save as-is
            combined_img = Image.fromarray(image_data_list[0][0])
        else:
            # Multiple images: combine them respecting original positions
            combined_img = combine_slide_images_with_positions(
                image_data_list, slide_width_emu, slide_height_emu
            )
        
        # Save the combined image
        name = f"slide{slide_idx:02d}.png"
        path = os.path.join(output_dir, name)
        combined_img.save(path)
        extracted.append(name)
    
    return extracted


def combine_slide_images_with_positions(image_data_list, slide_width_emu, slide_height_emu):
    """
    Combine multiple images from a slide into a single image, preserving their original positions.
    
    Parameters:
        image_data_list: List of tuples (img_array, left_emu, top_emu, width_emu, height_emu)
        slide_width_emu: Slide width in EMU units
        slide_height_emu: Slide height in EMU units
    
    Returns:
        PIL Image with all images placed at their original positions
    """
    if not image_data_list:
        raise ValueError("No images to combine")
    
    if len(image_data_list) == 1:
        return Image.fromarray(image_data_list[0][0])
    
    # Convert EMU to pixels
    # We'll use the actual image dimensions to determine the scale
    # First, find a reference scale based on the first image
    first_img, first_left, first_top, first_width_emu, first_height_emu = image_data_list[0]
    first_img_height, first_img_width = first_img.shape[:2]
    
    # Calculate pixels per EMU based on the first image
    # This assumes the image was placed at its natural size
    pixels_per_emu_height = first_img_height / first_height_emu if first_height_emu > 0 else 0
    pixels_per_emu_width = first_img_width / first_width_emu if first_width_emu > 0 else 0
    
    # Use average if both are valid, otherwise use the non-zero one
    if pixels_per_emu_height > 0 and pixels_per_emu_width > 0:
        pixels_per_emu = (pixels_per_emu_height + pixels_per_emu_width) / 2
    elif pixels_per_emu_height > 0:
        pixels_per_emu = pixels_per_emu_height
    elif pixels_per_emu_width > 0:
        pixels_per_emu = pixels_per_emu_width
    else:
        # Fallback: assume standard DPI (96 DPI = 914400 EMU per inch)
        # 1 inch = 96 pixels at 96 DPI, so 1 EMU = 96/914400 pixels
        pixels_per_emu = 96.0 / 914400.0
    
    # Calculate canvas size in pixels
    canvas_width = int(slide_width_emu * pixels_per_emu)
    canvas_height = int(slide_height_emu * pixels_per_emu)
    
    # Create canvas
    canvas = np.zeros((canvas_height, canvas_width, 3), dtype=np.uint8)
    
    # Place each image at its original position
    for img_array, left_emu, top_emu, width_emu, height_emu in image_data_list:
        # Convert positions to pixels
        left_px = int(left_emu * pixels_per_emu)
        top_px = int(top_emu * pixels_per_emu)
        width_px = int(width_emu * pixels_per_emu)
        height_px = int(height_emu * pixels_per_emu)
        
        # Resize image to match the shape's dimensions if needed
        img_height, img_width = img_array.shape[:2]
        if img_height != height_px or img_width != width_px:
            # Resize to match the shape dimensions
            img_pil = Image.fromarray(img_array)
            img_pil = img_pil.resize((width_px, height_px), Image.Resampling.LANCZOS)
            img_array = np.array(img_pil)
            img_height, img_width = img_array.shape[:2]
        
        # Calculate placement bounds
        x0 = max(0, left_px)
        y0 = max(0, top_px)
        x1 = min(canvas_width, left_px + img_width)
        y1 = min(canvas_height, top_px + img_height)
        
        # Calculate source bounds (in case image extends beyond canvas)
        src_x0 = max(0, -left_px)
        src_y0 = max(0, -top_px)
        src_x1 = src_x0 + (x1 - x0)
        src_y1 = src_y0 + (y1 - y0)
        
        # Place image on canvas
        if x1 > x0 and y1 > y0 and src_x1 > src_x0 and src_y1 > src_y0:
            new_img = img_array[src_y0:src_y1, src_x0:src_x1]
            existing = canvas[y0:y1, x0:x1]
            
            # Check if there's existing content at this exact position
            if np.any(existing > 0) and existing.shape == new_img.shape:
                # Same position and size: likely an overlay, use max blend (for fluorescence channels)
                canvas[y0:y1, x0:x1] = np.maximum(existing, new_img)
            else:
                # No overlap or different sizes: place the image (may overwrite if overlapping)
                # For non-overlapping images, this preserves their positions
                canvas[y0:y1, x0:x1] = new_img
    
    return Image.fromarray(canvas)

